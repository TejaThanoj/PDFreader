import pyttsx3
import PyPDF2
from tkinter.filedialog import askopenfilename
import os

# Add imports for PPTX and DOCX
from pptx import Presentation
from docx import Document

def read_pdf(book):
    pdfreader = PyPDF2.PdfReader(book)
    text = ""
    for page in pdfreader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def read_docx(book):
    doc = Document(book)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def read_pptx(book):
    prs = Presentation(book)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

book = askopenfilename()
ext = os.path.splitext(book)[1].lower()

if ext == ".pdf":
    text = read_pdf(book)
elif ext == ".docx":
    text = read_docx(book)
elif ext == ".pptx":
    text = read_pptx(book)
else:
    text = "Unsupported file type."

player = pyttsx3.init()
if text:
    player.say(text)
    player.runAndWait()